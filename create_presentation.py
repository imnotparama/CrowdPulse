from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define brand colors
    # AMD Red: #ED1C24 (237, 28, 36)
    # Background: Black (0, 0, 0)
    # Text: White (255, 255, 255)
    # Accent: Neon Blue (0, 255, 255)

    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

    def add_title_slide(prs, title_text, subtitle_text):
        slide_layout = prs.slide_layouts[0] # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        set_background(slide)

        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(237, 28, 36) # AMD Red
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(60)

        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)

    def add_content_slide(prs, title_text, bullet_points):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        set_background(slide)

        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(237, 28, 36) # AMD Red
        title.text_frame.paragraphs[0].font.bold = True

        body = slide.placeholders[1]
        tf = body.text_frame
        
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.size = Pt(20)
            p.space_after = Pt(10)

    # --- Slide 1: Title ---
    add_title_slide(prs, "CROWDPULSE", "AI-Powered Crowd Safety Intelligence\nDomain: AI for Smart Cities")

    # --- Slide 2: The Problem ---
    add_content_slide(prs, "The Problem: Urban Density Risks", [
        "Crowd crushes are preventable but deadly disasters (e.g., Itaewon 2022).",
        "Current monitoring relies on manual observation or post-event analysis.",
        "Lack of real-time data on 'Crowd Pressure' and 'Flow Dynamics'.",
        "Delayed response times during critical density surges.",
        "Infrastructure (stadiums, stations) often lacks digital safety intelligence."
    ])

    # --- Slide 3: The Solution ---
    add_content_slide(prs, "The Solution: CrowdPulse", [
        "A real-time AI surveillance layer for public safety.",
        "Instantly detects dangerous density, pressure surges, and flow stagnation.",
        "Transforms standard CCTV feeds into actionable safety metrics.",
        "Automated 'Evacuation Mode' alerts and PA announcements.",
        "Privacy-preserving analysis (hashed Wi-Fi probes, no facial recognition)."
    ])

    # --- Slide 4: Key Features ---
    add_content_slide(prs, "Key Features", [
        "📊 Real-time Density & Pressure Index: Quantifies physical compression risk.",
        "🚨 Stampede Risk Gauge: Predictive scoring (Safe → Critical).",
        "📡 Wi-Fi Probe Counting: Fusion of optical and RF sensor data.",
        "🗺️ Live Tactical Map: Dynamic evacuation routes and danger zones.",
        "📢 Automated PA System: Generates voice alerts based on crowd status.",
        "🆘 Emergency SOS: One-click site-wide evacuation protocol."
    ])

    # --- Slide 5: Tech Stack ---
    add_content_slide(prs, "Tech Stack", [
        "Computer Vision: YOLOv8 + OpenCV (Real-time detection & flow analysis)",
        "Backend: FastAPI (Python) + WebSocket (Live data streaming)",
        "Frontend: React + Vite + Recharts (Tactical Dashboard)",
        "IoT Simulation: ESP32 Wi-Fi Probe Sniffing (Crowd counting)",
        "Styling: TailwindCSS (High-contrast dark mode for command centers)"
    ])

    # --- Slide 6: Impact & Future ---
    add_content_slide(prs, "Impact: AI for Smart Cities", [
        "Prevents mass casualty events through early warning.",
        "Optimizes urban flow in transit hubs and festival grounds.",
        "Scalable to thousands of camera feeds via edge computing.",
        "Future: Integration with city-wide emergency dispatch (911/EMS).",
        "Future: Predictive crowd simulation for urban planning."
    ])

    prs.save('CrowdPulse_Presentation.pptx')
    print("Presentation saved successfully as CrowdPulse_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
